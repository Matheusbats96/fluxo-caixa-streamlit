import streamlit as st
import json
import os
import io
import pandas as pd
import plotly.express as px
import plotly.graph_objects as go
from pptx import Presentation
from pptx.util import Inches
from config import (
    RECEITA_MENSAL,
    MESES,
    IMPOSTOS_PCT,
    CONTINGENCIA_PCT,
    DESPESAS as DEFAULT_DESPESAS
)

# Arquivos de persistência
DESPESAS_FILE = 'despesas.json'
GROWTH_FILE = 'growth.json'
CONFIG_FILE = 'fluxo_config.json'

# Funções de JSON load/save
def load_json(path, default):
    if os.path.exists(path):
        with open(path, 'r', encoding='utf-8') as f:
            return json.load(f)
    return default

def save_json(data, path):
    with open(path, 'w', encoding='utf-8') as f:
        json.dump(data, f, ensure_ascii=False, indent=4)

# Inicialização de estado
if 'despesas' not in st.session_state:
    st.session_state.despesas = load_json(DESPESAS_FILE, DEFAULT_DESPESAS.copy())
if 'growth' not in st.session_state:
    default_growth = {k: 0.0 for k in DEFAULT_DESPESAS}
    st.session_state.growth = load_json(GROWTH_FILE, default_growth)
if 'config' not in st.session_state:
    default_cfg = {
        'receita_clube': RECEITA_MENSAL / 2,
        'receita_nao_membros': RECEITA_MENSAL / 2,
        'meses': MESES,
        'impostos_pct': IMPOSTOS_PCT,
        'contingencia_pct': CONTINGENCIA_PCT,
        'growth_receita': 0.0,
        'capex': 0.0
    }
    st.session_state.config = load_json(CONFIG_FILE, default_cfg)
# Garantir chaves
for key in ['receita_clube','receita_nao_membros','growth_receita','capex']:
    if key not in st.session_state.config:
        st.session_state.config[key] = 0.0 if key in ['growth_receita','capex'] else RECEITA_MENSAL/2
if 'df' not in st.session_state:
    st.session_state.df = None

# Título do app
st.title('Fluxo de Caixa Projetado')

# Seção de configurações iniciais
st.header('Configurações de Projeção')
cfg = st.session_state.config
new_receita_clube = st.number_input(
    'Receita Membros do Clube',
    value=float(cfg['receita_clube']),
    format='%.2f', step=0.01,
    key='input_receita_clube'
)
new_receita_nao = st.number_input(
    'Receita Não Membros',
    value=float(cfg['receita_nao_membros']),
    format='%.2f', step=0.01,
    key='input_receita_nao'
)
new_capex = st.number_input(
    'CAPEX (Implementação)',
    value=float(cfg['capex']),
    format='%.2f', step=0.01,
    key='input_capex'
)
new_meses = st.slider(
    'Meses de Projeção', min_value=1, max_value=36,
    value=cfg['meses'], key='slider_meses'
)
new_impostos = st.slider(
    'Taxa de Impostos (%)', min_value=0.0, max_value=0.2,
    value=cfg['impostos_pct'], format='%.2f',
    key='slider_impostos'
)
new_conting = st.slider(
    'Contingência (%)', min_value=0.0, max_value=0.2,
    value=cfg['contingencia_pct'], format='%.2f',
    key='slider_conting'
)
new_growth_r = st.number_input(
    'Crescimento mensal da Receita (%)',
    value=float(cfg.get('growth_receita',0.0)),
    format='%.2f', step=0.01,
    key='input_growth_receita'
)
# Atualizar config se alterar
if (
    new_receita_clube != cfg['receita_clube'] or
    new_receita_nao != cfg['receita_nao_membros'] or
    new_capex != cfg['capex'] or
    new_meses != cfg['meses'] or
    new_impostos != cfg['impostos_pct'] or
    new_conting != cfg['contingencia_pct'] or
    new_growth_r != cfg.get('growth_receita',0.0)
):
    st.session_state.config.update({
        'receita_clube': new_receita_clube,
        'receita_nao_membros': new_receita_nao,
        'capex': new_capex,
        'meses': new_meses,
        'impostos_pct': new_impostos,
        'contingencia_pct': new_conting,
        'growth_receita': new_growth_r
    })
    save_json(st.session_state.config, CONFIG_FILE)

# Seção de despesas OPEX
st.header('Despesas/Custos: ')
remove_list = []
for nome, val in st.session_state.despesas.items():
    c1, c2, c3 = st.columns([3,2,1])
    new_val = c1.number_input(
        nome,
        value=float(val),
        format='%.2f', step=0.01,
        key=f'val_{nome}'
    )
    new_g = c2.number_input(
        f'Crescimento anual {nome} (%)',
        value=float(st.session_state.growth.get(nome,0.0)),
        format='%.2f', step=0.01,
        key=f'gr_{nome}'
    )
    if c3.button('Excluir', key=f'del_{nome}'):
        remove_list.append(nome)
    else:
        st.session_state.despesas[nome] = new_val
        st.session_state.growth[nome] = new_g
for r in remove_list:
    st.session_state.despesas.pop(r, None)
    st.session_state.growth.pop(r, None)
save_json(st.session_state.despesas, DESPESAS_FILE)
save_json(st.session_state.growth, GROWTH_FILE)

# Adicionar novo item
st.header('Adicionar Novo Item')
with st.form('add_form'):
    nm = st.text_input('Nome do item')
    nv = st.number_input('Valor mensal', min_value=0.0, format='%.2f', step=0.01)
    ng = st.number_input('Crescimento anual (%)', min_value=0.0, format='%.2f', step=0.01)
    if st.form_submit_button('Adicionar'):
        if nm.strip() and nm not in st.session_state.despesas:
            st.session_state.despesas[nm] = nv
            st.session_state.growth[nm] = ng
            save_json(st.session_state.despesas, DESPESAS_FILE)
            save_json(st.session_state.growth, GROWTH_FILE)
            st.success(f'Item \"{nm}\" adicionado.')
            st.rerun()
        else:
            st.error('Nome inválido ou duplicado.')

# Geração do fluxo de caixa
if st.button('Gerar Fluxo de Caixa'):
    cfg = st.session_state.config
    dates = [pd.to_datetime('2025-08-01') + pd.DateOffset(months=i) for i in range(cfg['meses'])]
    dfc = pd.DataFrame(index=dates)
    dfc.index.name = 'Mês'

    # Fator de growth mensal
    gm = (1 + cfg.get('growth_receita',0.0))**(1/12) - 1

    # Receitas segmentadas
    rc = cfg['receita_clube']
    rn = cfg['receita_nao_membros']
    dfc['Receita Clube'] = [rc * (1+gm)**i for i in range(cfg['meses'])]
    dfc['Receita Não Membros'] = [rn * (1+gm)**i for i in range(cfg['meses'])]
    dfc['Receita'] = dfc['Receita Clube'] + dfc['Receita Não Membros']

    # Despesa de Capex pontual
    dfc['Capex'] = [cfg.get('capex',0.0)] + [0]*(cfg['meses']-1)

    # Despesas OPEX com growth
    for nome,val in st.session_state.despesas.items():
        gd = (1 + st.session_state.growth.get(nome,0.0))**(1/12) - 1
        dfc[nome] = [val * (1+gd)**i for i in range(cfg['meses'])]

    # --- Início dos Cálculos ---
    opex_cols = list(st.session_state.despesas.keys())
    dfc['OPEX Total'] = dfc[opex_cols].sum(axis=1)
    dfc['Contingência'] = dfc['OPEX Total'] * cfg['contingencia_pct']

    impostos_col_name = f'Impostos ({int(cfg["impostos_pct"]*100)}%)'
    dfc[impostos_col_name] = dfc['Receita'] * cfg['impostos_pct']

    despesas_totais_cols = ['OPEX Total', 'Contingência', impostos_col_name, 'Capex']
    dfc['Despesas Totais'] = dfc[despesas_totais_cols].sum(axis=1)

    dfc['Lucro Operacional'] = dfc['Receita'] - dfc['Despesas Totais']
    
    # Esta coluna é necessária para o cálculo do payback
    dfc['Fluxo Acumulado'] = dfc['Lucro Operacional'].cumsum()
    # --- Fim dos Cálculos ---

    st.session_state.df = dfc
    st.rerun() # Adicionado para garantir que a exibição seja atualizada

# Exibição e análise
if st.session_state.df is not None:
    df = st.session_state.df
    st.subheader('Fluxo de Caixa (Tabela)')
    st.dataframe(df.style.format("{:,.2f}"))

    # KPIs
    total_receita_kpi = df['Receita'].sum()
    avg_m = (df['Lucro Operacional'].sum() / total_receita_kpi) if total_receita_kpi > 0 else 0

    tot_r = df['Receita'].sum()
    tot_d = df['Despesas Totais'].sum()
    tot_p = df['Lucro Operacional'].sum()
    mmx = df['Lucro Operacional'].idxmax().strftime('%b %Y')
    mmn = df['Lucro Operacional'].idxmin().strftime('%b %Y')
    
    # --- INÍCIO DA CORREÇÃO ---
    # O cálculo do payback foi movido para cá
    payback_date = df[df['Fluxo Acumulado'] >= 0].index
    payback = payback_date[0].strftime('%b %Y') if len(payback_date) > 0 else 'Não alcançado'
    # --- FIM DA CORREÇÃO ---

    st.markdown(f"**Margem Média:**<br><span style='font-size:24px'>{avg_m:.2%}</span>", unsafe_allow_html=True)
    st.markdown(f"**Receita Total:**<br><span style='font-size:24px'>R$ {tot_r:,.2f}</span>", unsafe_allow_html=True)
    st.markdown(f"**Despesas Totais:**<br><span style='font-size:24px'>R$ {tot_d:,.2f}</span>", unsafe_allow_html=True)
    st.markdown(f"**Lucro Total:**<br><span style='font-size:24px'>R$ {tot_p:,.2f}</span>", unsafe_allow_html=True)
    st.markdown(f"**Payback (primeiro mês ≥0):**<br><span style='font-size:24px'>{payback}</span>", unsafe_allow_html=True)
    st.markdown(f"**Mês Maior Lucro:**<br><span style='font-size:24px'>{mmx}</span>", unsafe_allow_html=True)
    st.markdown(f"**Mês Menor Lucro:**<br><span style='font-size:24px'>{mmn}</span>", unsafe_allow_html=True)

    dff = df.reset_index().rename(columns={'index':'Mês'})

    figl = px.line(
        dff,
        x='Mês',
        y=['Receita', 'Despesas Totais', 'Fluxo Acumulado'],
        title='Receita, Despesas Totais e Fluxo de Caixa Acumulado'
    )
    figl.update_traces(selector={'name': 'Despesas Totais'}, y=-dff['Despesas Totais'])
    st.plotly_chart(figl, use_container_width=True)

    # --- WATERFALL (ESTILO REVERTIDO, ORDENAÇÃO MANTIDA) ---
    impostos_col_name = [col for col in df.columns if 'Impostos' in col][0]
    total_rev = df['Receita'].sum()
    total_opex = df['OPEX Total'].sum()
    total_imp = df[impostos_col_name].sum()
    total_cont = df['Contingência'].sum()
    total_cap = df['Capex'].sum()
    total_lucro = df['Lucro Operacional'].sum()

    receita_item = {'category': 'Receita Total', 'value': total_rev, 'measure': 'absolute'}
    lucro_item = {'category': 'Lucro Total', 'value': total_lucro, 'measure': 'total'}
    despesa_items = [
        {'category': 'OPEX', 'value': -total_opex, 'measure': 'relative'},
        {'category': 'Impostos', 'value': -total_imp, 'measure': 'relative'},
        {'category': 'Contingência', 'value': -total_cont, 'measure': 'relative'},
        {'category': 'Capex', 'value': -total_cap, 'measure': 'relative'},
    ]
    despesa_items = [item for item in despesa_items if item['value'] != 0]
    sorted_despesa_items = sorted(despesa_items, key=lambda x: x['value'])

    wdata = [receita_item] + sorted_despesa_items + [lucro_item]
    wdf = pd.DataFrame(wdata)

    figw = go.Figure(go.Waterfall(
        x=wdf['category'],
        y=wdf['value'],
        measure=wdf['measure'],
        text=wdf['value'].map(lambda v: f"R$ {v:,.2f}"),
        textposition='outside',
        connector={'line': {'color': '#7f8c8d', 'width': 1}},
        increasing={'marker': {'color': '#27ae60'}},
        decreasing={'marker': {'color': '#c0392b'}},
        totals={'marker': {'color': '#2980b9'}}
    ))
    figw.update_layout(
        title='Receita/Lucro',
        yaxis={'title': 'Valor (R$)', 'tickprefix': 'R$ '},
        xaxis={'tickangle': -45},
        plot_bgcolor='aliceblue',
        font={'family': 'Arial', 'size': 7, 'color': '#34495e'},
        margin={'l': 40, 'r': 40, 't': 80, 'b': 120}
    )
    st.plotly_chart(figw, use_container_width=True)


    # --- PIE CHART (LÓGICA CORRETA, ESTILO DE TEXTO REVERTIDO) ---
    opex_cols = list(st.session_state.despesas.keys())
    opex_totals_by_item = df[opex_cols].sum()
    opex_totals_by_item = opex_totals_by_item[opex_totals_by_item > 0]

    if not opex_totals_by_item.empty:
        pdp = pd.DataFrame(opex_totals_by_item).reset_index()
        pdp.columns = ['Despesa', 'Valor']
        
        figp = px.pie(pdp,
            names='Despesa',
            values='Valor',
            title='Composição Percentual das Despesas OPEX'
        )
        st.plotly_chart(figp, use_container_width=True)


    # Export Slides
    if st.button('Exportar Slides'):
        prs = Presentation()
        s0 = prs.slides.add_slide(prs.slide_layouts[5])
        s0.shapes.title.text = 'Resumo do Fluxo de Caixa'
        tb = s0.shapes.add_textbox(Inches(1),Inches(1.5),Inches(8),Inches(2.5))
        tf = tb.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = f'Receita Total: R$ {tot_r:,.2f}'
        p = tf.add_paragraph()
        p.text = f'Despesas Totais: R$ {tot_d:,.2f}'
        p = tf.add_paragraph()
        p.text = f'Lucro Total: R$ {tot_p:,.2f}'
        p = tf.add_paragraph()
        p.text = f'Payback: {payback}'

        img_l = io.BytesIO(figl.to_image(format='png'))
        img_w = io.BytesIO(figw.to_image(format='png'))

        sl1 = prs.slides.add_slide(prs.slide_layouts[5])
        sl1.shapes.title.text = figl.layout.title.text
        sl1.shapes.add_picture(img_l, Inches(0.5), Inches(1.5), width=Inches(9))

        sl2 = prs.slides.add_slide(prs.slide_layouts[5])
        sl2.shapes.title.text = figw.layout.title.text
        sl2.shapes.add_picture(img_w, Inches(0.5), Inches(1.5), width=Inches(9))

        if 'figp' in locals():
            img_p = io.BytesIO(figp.to_image(format='png'))
            sl3 = prs.slides.add_slide(prs.slide_layouts[5])
            sl3.shapes.title.text = figp.layout.title.text
            sl3.shapes.add_picture(img_p, Inches(1), Inches(1.5), height=Inches(5))

        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)

        st.download_button(
            'Download Slides',
            data=buf,
            file_name='fluxo_slides.pptx',
            mime='application/vnd.openxmlformats-officedocument.presentationml.presentation'
        )
